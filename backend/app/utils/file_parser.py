"""
文件解析工具
支持 Word (docx/doc)、Excel (xlsx/xls/csv)、PPT (pptx)、PDF、网页 (html)、数据 (json/rtf/log)、图片截图 (png/jpg/webp) 及纯文本 (txt/md) 的深度文本与结构提取
"""

import os
import io
import csv
import json
import base64
from pathlib import Path
from typing import List, Optional


def _read_text_with_fallback(file_path: str) -> str:
    """
    读取文本文件，UTF-8失败时自动探测编码。
    采用多级回退策略：
    1. UTF-8
    2. charset_normalizer
    3. chardet
    4. UTF-8 + errors='replace'
    """
    data = Path(file_path).read_bytes()
    
    # 1. 尝试 UTF-8
    try:
        return data.decode('utf-8')
    except UnicodeDecodeError:
        pass
    
    # 2. 尝试 charset_normalizer
    encoding = None
    try:
        from charset_normalizer import from_bytes
        best = from_bytes(data).best()
        if best and best.encoding:
            encoding = best.encoding
    except Exception:
        pass
    
    # 3. 回退到 chardet
    if not encoding:
        try:
            import chardet
            result = chardet.detect(data)
            encoding = result.get('encoding') if result else None
        except Exception:
            pass
    
    # 4. 最终兜底
    if not encoding:
        encoding = 'utf-8'
    
    return data.decode(encoding, errors='replace')


class FileParser:
    """多格式个人资料解析器"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv', '.pptx',
        '.txt', '.md', '.markdown', '.html', '.htm', '.json', '.log', '.rtf',
        '.png', '.jpg', '.jpeg', '.webp'
    }
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """检查文件是否为支持的格式"""
        suffix = Path(file_path).suffix.lower()
        return suffix in cls.SUPPORTED_EXTENSIONS
    
    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """从文件中提取文本"""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        suffix = path.suffix.lower()
        
        if suffix not in cls.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的文件格式: {suffix}")
        
        if suffix == '.pdf':
            return cls._extract_from_pdf(file_path)
        elif suffix == '.docx':
            return cls._extract_from_docx(file_path)
        elif suffix == '.doc':
            return cls._extract_from_doc(file_path)
        elif suffix in {'.xlsx', '.xls'}:
            return cls._extract_from_xlsx(file_path)
        elif suffix == '.csv':
            return cls._extract_from_csv(file_path)
        elif suffix == '.pptx':
            return cls._extract_from_pptx(file_path)
        elif suffix in {'.html', '.htm'}:
            return cls._extract_from_html(file_path)
        elif suffix == '.json':
            return cls._extract_from_json(file_path)
        elif suffix in {'.png', '.jpg', '.jpeg', '.webp'}:
            return cls._extract_from_image(file_path)
        elif suffix in {'.md', '.markdown'}:
            return cls._extract_from_md(file_path)
        elif suffix in {'.txt', '.log', '.rtf'}:
            return cls._extract_from_txt(file_path)
        
        return _read_text_with_fallback(file_path)
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """从 PDF 提取文本"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("需要安装PyMuPDF: pip install PyMuPDF")
        
        text_parts = []
        with fitz.open(file_path) as doc:
            for page in doc:
                text = page.get_text()
                if text.strip():
                    text_parts.append(text)
        
        return "\n\n".join(text_parts)
    
    @staticmethod
    def _extract_from_docx(file_path: str) -> str:
        """从 Word (.docx) 提取文本与表格内容"""
        try:
            import docx
        except ImportError:
            raise ImportError("需要安装python-docx: pip install python-docx")
        
        doc = docx.Document(file_path)
        parts = []
        
        # 1. 段落正文
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                parts.append(t)
        
        # 2. 表格数据（如简历中的经历表、成绩单、技能矩阵）
        for table in doc.tables:
            table_lines = []
            for row in table.rows:
                cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                # 去除全空行
                if any(cells):
                    table_lines.append(" | ".join(cells))
            if table_lines:
                parts.append("\n".join(table_lines))
        
        return "\n\n".join(parts)
    
    @staticmethod
    def _extract_from_doc(file_path: str) -> str:
        """从旧版 Word (.doc) 尝试提取纯文本内容"""
        # 针对旧版 doc 二进制流，先尝试 docx 解析，失败则回退到多编码流提取
        try:
            return FileParser._extract_from_docx(file_path)
        except Exception:
            pass
        return _read_text_with_fallback(file_path)
    
    @staticmethod
    def _extract_from_xlsx(file_path: str) -> str:
        """从 Excel (.xlsx/.xls) 提取各 Sheet 的表格与时间线记录"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("需要安装openpyxl: pip install openpyxl")
        
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        except Exception:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
        sheet_texts = []
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                if not row or not any(row):
                    continue
                cleaned = [str(c).strip() if c is not None else "" for c in row]
                if any(cleaned):
                    rows_text.append(" | ".join(cleaned))
            if rows_text:
                sheet_texts.append(f"### 工作表: {sheetname}\n" + "\n".join(rows_text))
        
        wb.close()
        return "\n\n".join(sheet_texts)
    
    @staticmethod
    def _extract_from_csv(file_path: str) -> str:
        """从 CSV 提取结构化记录"""
        raw_text = _read_text_with_fallback(file_path)
        try:
            reader = csv.reader(io.StringIO(raw_text))
            lines = []
            for row in reader:
                if any(row):
                    lines.append(" | ".join([c.strip() for c in row]))
            if lines:
                return "\n".join(lines)
        except Exception:
            pass
        return raw_text
    
    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        """从 PowerPoint (.pptx) 提取幻灯片标题、文本框与演讲备注"""
        try:
            import pptx
        except ImportError:
            raise ImportError("需要安装python-pptx: pip install python-pptx")
        
        prs = pptx.Presentation(file_path)
        slide_texts = []
        for idx, slide in enumerate(prs.slides, 1):
            cur_slide = [f"### 幻灯片 {idx}"]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        cur_slide.append(text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    cur_slide.append(f"[演讲备注]: {notes}")
            if len(cur_slide) > 1:
                slide_texts.append("\n".join(cur_slide))
        
        return "\n\n".join(slide_texts)
    
    @staticmethod
    def _extract_from_html(file_path: str) -> str:
        """从 HTML 网页或笔记导出包清洗并提取正文"""
        try:
            from bs4 import BeautifulSoup
            raw_html = _read_text_with_fallback(file_path)
            soup = BeautifulSoup(raw_html, 'html.parser')
            for script in soup(["script", "style", "meta", "noscript", "link"]):
                script.extract()
            text = soup.get_text(separator="\n")
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            return "\n".join(lines)
        except Exception:
            return _read_text_with_fallback(file_path)
    
    @staticmethod
    def _extract_from_json(file_path: str) -> str:
        """从 JSON 格式数据（如聊天记录导出、问卷数据）解析为可读文本"""
        raw_text = _read_text_with_fallback(file_path)
        try:
            data = json.loads(raw_text)
            if isinstance(data, list):
                formatted = []
                for item in data:
                    if isinstance(item, dict):
                        sender = item.get('sender') or item.get('name') or item.get('role') or item.get('user') or '记录'
                        msg = item.get('content') or item.get('message') or item.get('text') or str(item)
                        time_str = item.get('time') or item.get('date') or item.get('timestamp') or ''
                        prefix = f"[{time_str}] " if time_str else ""
                        formatted.append(f"{prefix}{sender}: {msg}")
                    else:
                        formatted.append(str(item))
                return "\n".join(formatted)
            elif isinstance(data, dict):
                return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception:
            pass
        return raw_text
    
    @staticmethod
    def _extract_from_image(file_path: str) -> str:
        """通过大模型多模态 Vision 接口提取图片文字内容（聊天截图/手写日记/总结图）"""
        try:
            from app.utils.llm_client import LLMClient
            from app.config import Config
            
            with open(file_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            
            ext = Path(file_path).suffix.lower().lstrip('.')
            if ext == 'jpg':
                ext = 'jpeg'
            mime_type = f"image/{ext}"
            
            prompt = (
                "你是一个精准的 OCR 与生活资料信息提取器。请完整、客观地提取这张图片中的所有文字内容、聊天记录、日记随笔或表格信息。"
                "如果是聊天记录截图，请按'发言人: 内容'的格式还原；"
                "如果是手写日记或文章照片，请保持原始段落还原；"
                "如果是表格或表单，请转为清晰的文本列表；"
                "只输出提取到的文字内容，不要输出任何额外的说明或寒暄。"
            )
            
            client = LLMClient()
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{encoded_string}"
                            }
                        }
                    ]
                }
            ]
            response = client.chat(messages=messages, temperature=0.1)
            return response.strip()
        except Exception as e:
            # 优雅降级：提示图片已收录但需用户核验
            filename = Path(file_path).name
            return f"[图片资料: {filename} (OCR提取回退: {str(e)})]"
    
    @staticmethod
    def _extract_from_md(file_path: str) -> str:
        """从 Markdown 提取文本"""
        return _read_text_with_fallback(file_path)
    
    @staticmethod
    def _extract_from_txt(file_path: str) -> str:
        """从 TXT/LOG/RTF 提取文本"""
        return _read_text_with_fallback(file_path)
    
    @classmethod
    def extract_from_multiple(cls, file_paths: List[str]) -> str:
        """从多个文件提取文本并合并"""
        all_texts = []
        for i, file_path in enumerate(file_paths, 1):
            try:
                text = cls.extract_text(file_path)
                filename = Path(file_path).name
                all_texts.append(f"=== 文档 {i}: {filename} ===\n{text}")
            except Exception as e:
                all_texts.append(f"=== 文档 {i}: {file_path} (提取失败: {str(e)}) ===")
        return "\n\n".join(all_texts)


def split_text_into_chunks(
    text: str, 
    chunk_size: int = 500, 
    overlap: int = 50
) -> List[str]:
    """将文本分割成小块"""
    if len(text) <= chunk_size:
        return [text] if text.strip() else []
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            for sep in ['。', '！', '？', '.\n', '!\n', '?\n', '\n\n', '. ', '! ', '? ']:
                last_sep = text[start:end].rfind(sep)
                if last_sep != -1 and last_sep > chunk_size * 0.3:
                    end = start + last_sep + len(sep)
                    break
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        start = end - overlap if end < len(text) else len(text)
    
    return chunks
