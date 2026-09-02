# -*- coding: utf-8 -*-
"""
Tests for multi-format FileParser
"""

import os
import tempfile
import pytest
from app.utils.file_parser import FileParser, split_text_into_chunks


def test_is_supported():
    assert FileParser.is_supported("resume.docx") is True
    assert FileParser.is_supported("financials.xlsx") is True
    assert FileParser.is_supported("timeline.csv") is True
    assert FileParser.is_supported("summary.pptx") is True
    assert FileParser.is_supported("chat.html") is True
    assert FileParser.is_supported("messages.json") is True
    assert FileParser.is_supported("diary.png") is True
    assert FileParser.is_supported("photo.jpg") is True
    assert FileParser.is_supported("doc.txt") is True
    assert FileParser.is_supported("doc.md") is True
    assert FileParser.is_supported("archive.zip") is False


def test_extract_from_docx():
    import docx
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        path = f.name
    try:
        doc = docx.Document()
        doc.add_paragraph("张三的个人简历与经历总结")
        doc.add_paragraph("曾在科技公司担任核心工程师，负责系统架构。")
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "年份"
        table.rows[0].cells[1].text = "事件"
        table.rows[1].cells[0].text = "2024"
        table.rows[1].cells[1].text = "决定独立创业"
        doc.save(path)

        text = FileParser.extract_text(path)
        assert "张三的个人简历" in text
        assert "核心工程师" in text
        assert "2024 | 决定独立创业" in text
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_extract_from_xlsx():
    import openpyxl
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        path = f.name
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "财务与储蓄流水"
        ws.append(["日期", "项目", "金额", "结余现金流月数"])
        ws.append(["2025-01", "存款储备", "150000", "12个月"])
        wb.save(path)

        text = FileParser.extract_text(path)
        assert "工作表: 财务与储蓄流水" in text
        assert "日期 | 项目 | 金额 | 结余现金流月数" in text
        assert "150000 | 12个月" in text
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_extract_from_csv():
    with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", encoding="utf-8", delete=False) as f:
        f.write("时间,地点,心情,事件\n2025-02-01,北京,焦虑,与朋友讨论人生规划\n")
        path = f.name
    try:
        text = FileParser.extract_text(path)
        assert "时间 | 地点 | 心情 | 事件" in text
        assert "2025-02-01 | 北京 | 焦虑 | 与朋友讨论人生规划" in text
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_extract_from_pptx():
    import pptx
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    try:
        prs = pptx.Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        tx_box = slide.shapes.add_textbox(100, 100, 400, 200)
        tf = tx_box.text_frame
        tf.text = "个人十年愿景与核心战略汇报"
        p = tf.add_paragraph()
        p.text = "核心目标：在混沌中建立自己的选择权体系。"
        prs.save(path)

        text = FileParser.extract_text(path)
        assert "幻灯片 1" in text
        assert "个人十年愿景" in text
        assert "选择权体系" in text
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_extract_from_html():
    with tempfile.NamedTemporaryFile(suffix=".html", mode="w", encoding="utf-8", delete=False) as f:
        f.write("<html><head><title>日记</title><style>body{color:red;}</style></head><body><h1>2025年感悟</h1><p>今天读了特德姜的书，对未来有了新的思考。</p></body></html>")
        path = f.name
    try:
        text = FileParser.extract_text(path)
        assert "2025年感悟" in text
        assert "今天读了特德姜的书" in text
        assert "body{color:red;}" not in text
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_extract_from_json():
    with tempfile.NamedTemporaryFile(suffix=".json", mode="w", encoding="utf-8", delete=False) as f:
        f.write('[{"time": "14:20", "sender": "好友小李", "content": "不要慌，不管选哪条路我们都一起想办法。"}]')
        path = f.name
    try:
        text = FileParser.extract_text(path)
        assert "好友小李" in text
        assert "不要慌" in text
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_split_text_into_chunks():
    long_text = "这是第一句话。这是第二句话！这是第三句话？这是第四句话。" * 30
    chunks = split_text_into_chunks(long_text, chunk_size=100, overlap=20)
    assert len(chunks) > 1
    for c in chunks:
        assert len(c) > 0
